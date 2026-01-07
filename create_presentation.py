from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
import fitz
import os

# PDF to PNG conversion
def pdf_to_png(pdf_path, output_path, dpi=150):
    if not os.path.exists(pdf_path):
        print(f"Warning: {pdf_path} not found")
        return None
    doc = fitz.open(pdf_path)
    page = doc[0]
    mat = fitz.Matrix(dpi/72, dpi/72)
    pix = page.get_pixmap(matrix=mat)
    pix.save(output_path)
    doc.close()
    return output_path

# Convert all PDFs to PNGs
figures_dir = "/Users/umutturklay/dev/UE/pattern_recog/Figures_Tables"
png_dir = "/Users/umutturklay/dev/UE/pattern_recog/Figures_PNG"
os.makedirs(png_dir, exist_ok=True)

pdf_files = {
    "RQ1_Fig1": f"{figures_dir}/RQ1/RQ1_Fig1.pdf",
    "RQ1_Fig2": f"{figures_dir}/RQ1/RQ1_Fig2.pdf",
    "RQ1_Fig3": f"{figures_dir}/RQ1/RQ1_Fig3.pdf",
    "RQ1_Fig4": f"{figures_dir}/RQ1/RQ1_Fig4.pdf",
    "RQ2_Fig1": f"{figures_dir}/RQ2/RQ2_Fig1.pdf",
    "RQ3_Fig1": f"{figures_dir}/RQ3/RQ3_Fig1.pdf",
    "RQ4_Fig1": f"{figures_dir}/RQ4/RQ4_Fig1.pdf",
    "RQ5_Fig1": f"{figures_dir}/RQ5/RQ5_Fig1.pdf",
}

png_files = {}
for name, pdf_path in pdf_files.items():
    png_path = f"{png_dir}/{name}.png"
    result = pdf_to_png(pdf_path, png_path)
    if result:
        png_files[name] = png_path
        print(f"Converted: {name}")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme - Medical/Health theme
PRIMARY_BLUE = RGBColor(0, 102, 153)      # Deep medical blue
SECONDARY_BLUE = RGBColor(0, 153, 204)    # Lighter blue
ACCENT_TEAL = RGBColor(0, 166, 147)       # Teal/turquoise
LIGHT_BG = RGBColor(240, 248, 255)        # Alice blue (very light)
DARK_TEXT = RGBColor(33, 37, 41)          # Dark gray
WHITE = RGBColor(255, 255, 255)
HEADER_HEIGHT = Inches(1.1)

def add_slide_design(slide, slide_num=None, total_slides=14):
    """Add medical-themed design elements to slide"""
    # Top header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, HEADER_HEIGHT)
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_BLUE
    header.line.fill.background()

    # Accent line under header
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, HEADER_HEIGHT, prs.slide_width, Inches(0.05))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_TEAL
    accent.line.fill.background()

    # Bottom footer bar
    footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.1), prs.slide_width, Inches(0.4))
    footer.fill.solid()
    footer.fill.fore_color.rgb = PRIMARY_BLUE
    footer.line.fill.background()

    # Footer text - left
    footer_left = slide.shapes.add_textbox(Inches(0.3), Inches(7.15), Inches(6), Inches(0.3))
    tf = footer_left.text_frame
    p = tf.paragraphs[0]
    p.text = "Pattern Recognition | Group 8 | Brain Tumor Diagnosis"
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE

    # Footer text - right (slide number)
    if slide_num:
        footer_right = slide.shapes.add_textbox(Inches(11.5), Inches(7.15), Inches(1.5), Inches(0.3))
        tf = footer_right.text_frame
        p = tf.paragraphs[0]
        p.text = f"{slide_num} / {total_slides}"
        p.font.size = Pt(10)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.RIGHT

    # Decorative medical cross (subtle)
    cross_size = Inches(0.15)
    cross_x = Inches(12.8)
    cross_y = Inches(0.45)

    # Vertical part
    v_cross = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cross_x, cross_y - cross_size/2, Inches(0.05), cross_size)
    v_cross.fill.solid()
    v_cross.fill.fore_color.rgb = WHITE
    v_cross.line.fill.background()

    # Horizontal part
    h_cross = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, cross_x - cross_size/2 + Inches(0.025), cross_y - Inches(0.025), cross_size, Inches(0.05))
    h_cross.fill.solid()
    h_cross.fill.fore_color.rgb = WHITE
    h_cross.line.fill.background()

def add_title_slide(prs, title, subtitle=None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Full background gradient effect with shapes
    bg1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = PRIMARY_BLUE
    bg1.line.fill.background()

    # Decorative circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-2), Inches(8), Inches(8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = SECONDARY_BLUE
    circle.line.fill.background()

    # Another decorative element
    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-3), Inches(4), Inches(6), Inches(6))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = ACCENT_TEAL
    circle2.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(2.5))
        tf = sub_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(22)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Bottom accent line
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(3.7), Inches(5.333), Inches(0.05))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_TEAL
    accent.line.fill.background()

    return slide

def add_content_slide(prs, title, bullet_points, slide_num=None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_slide_design(slide, slide_num)

    # Title in header
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(11), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Content area with slight background
    content_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.3), Inches(12.733), Inches(5.6))
    content_bg.fill.solid()
    content_bg.fill.fore_color.rgb = RGBColor(250, 252, 255)
    content_bg.line.color.rgb = RGBColor(200, 220, 240)

    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.133), Inches(5.3))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if point.startswith("   "):
            p.text = f"    {point.strip()}"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(80, 80, 80)
        elif point == "":
            p.text = ""
            p.font.size = Pt(12)
        else:
            p.text = f"● {point}"
            p.font.size = Pt(22)
            p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(8)

    return slide

def add_table_slide(prs, title, headers, rows, note=None, slide_num=None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_slide_design(slide, slide_num)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(11), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    cols = len(headers)
    row_count = len(rows) + 1
    table_width = Inches(12.5)
    cell_height = 0.5
    table_height = Inches(cell_height * row_count)
    left = Inches(0.417)
    top = Inches(1.4)

    table = slide.shapes.add_table(row_count, cols, left, top, table_width, table_height).table

    # Header row styling
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_BLUE

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_data)
            cell.text_frame.paragraphs[0].font.size = Pt(13)
            cell.text_frame.paragraphs[0].font.color.rgb = DARK_TEXT
            # Alternating row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 248, 255)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    if note:
        note_top = top + table_height + Inches(0.3)
        note_box = slide.shapes.add_textbox(Inches(0.5), note_top, Inches(12.333), Inches(1.5))
        tf = note_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = note
        p.font.size = Pt(16)
        p.font.italic = True
        p.font.color.rgb = RGBColor(80, 80, 80)

    return slide

def add_figure_slide(prs, title, image_path, bullets=None, slide_num=None, fullsize=False):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_slide_design(slide, slide_num)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(11), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if fullsize:
        # Full-size figure mode - figure takes almost entire slide
        fig_height = Inches(5.7)
        fig_top = Inches(1.2)
        fig_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), fig_top, Inches(12.733), fig_height)
        fig_bg.fill.solid()
        fig_bg.fill.fore_color.rgb = WHITE
        fig_bg.line.color.rgb = RGBColor(200, 220, 240)
        fig_bg.line.width = Pt(1)

        if image_path and os.path.exists(image_path):
            # Use width for landscape figures, centered horizontally
            pic = slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.3), width=Inches(12.333))
            # Center the image if it's not as wide as expected
            pic_left = (prs.slide_width - pic.width) / 2
            pic.left = int(pic_left)
    else:
        # Standard figure mode with bullets
        fig_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.25), Inches(12.733), Inches(4.5))
        fig_bg.fill.solid()
        fig_bg.fill.fore_color.rgb = WHITE
        fig_bg.line.color.rgb = RGBColor(200, 220, 240)
        fig_bg.line.width = Pt(1)

        if image_path and os.path.exists(image_path):
            slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.35), width=Inches(12.333), height=Inches(4.3))
        else:
            # Placeholder
            text_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.733), Inches(0.8))
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            p.text = "[INSERT FIGURE]"
            p.font.size = Pt(18)
            p.font.italic = True
            p.font.color.rgb = RGBColor(150, 150, 150)
            p.alignment = PP_ALIGN.CENTER

        if bullets:
            # Bullet area background
            bullet_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(5.85), Inches(12.733), Inches(1.15))
            bullet_bg.fill.solid()
            bullet_bg.fill.fore_color.rgb = RGBColor(240, 248, 255)
            bullet_bg.line.color.rgb = RGBColor(200, 220, 240)

            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.95), Inches(12.333), Inches(1))
            tf = content_box.text_frame
            tf.word_wrap = True
            for i, point in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"● {point}"
                p.font.size = Pt(14)
                p.font.color.rgb = DARK_TEXT

    return slide

def add_rq_result_slide(prs, rq_num, rq_title, image_path, table_headers, table_rows, observations, slide_num=None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_slide_design(slide, slide_num)

    # Title with RQ badge
    # RQ badge
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.3), Inches(0.2), Inches(0.7), Inches(0.7))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT_TEAL
    badge.line.fill.background()

    badge_text = slide.shapes.add_textbox(Inches(0.3), Inches(0.35), Inches(0.7), Inches(0.4))
    tf = badge_text.text_frame
    p = tf.paragraphs[0]
    p.text = f"RQ{rq_num}"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Title
    title_box = slide.shapes.add_textbox(Inches(1.1), Inches(0.25), Inches(11), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = rq_title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Figure area (left)
    fig_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.2), Inches(1.2), Inches(6.2), Inches(3.7))
    fig_bg.fill.solid()
    fig_bg.fill.fore_color.rgb = WHITE
    fig_bg.line.color.rgb = RGBColor(200, 220, 240)

    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(0.3), Inches(1.3), width=Inches(6), height=Inches(3.5))
    else:
        fig_text = slide.shapes.add_textbox(Inches(0.3), Inches(2.8), Inches(6), Inches(0.5))
        tf = fig_text.text_frame
        p = tf.paragraphs[0]
        p.text = "[INSERT FIGURE]"
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = RGBColor(150, 150, 150)
        p.alignment = PP_ALIGN.CENTER

    # Table (right)
    if table_headers and table_rows:
        cols = len(table_headers)
        row_count = len(table_rows) + 1
        table_width = Inches(6.6)
        table_height = Inches(min(3.5, 0.45 * row_count))

        table = slide.shapes.add_table(row_count, cols, Inches(6.5), Inches(1.2), table_width, table_height).table

        for i, header in enumerate(table_headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = PRIMARY_BLUE

        for row_idx, row_data in enumerate(table_rows):
            for col_idx, cell_data in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(cell_data)
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(240, 248, 255)

    # Key Observations box
    obs_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.2), Inches(5.0), Inches(12.9), Inches(2))
    obs_bg.fill.solid()
    obs_bg.fill.fore_color.rgb = RGBColor(232, 245, 233)  # Light green
    obs_bg.line.color.rgb = ACCENT_TEAL

    obs_box = slide.shapes.add_textbox(Inches(0.4), Inches(5.1), Inches(12.5), Inches(1.9))
    tf = obs_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Key Observations"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE

    for obs in observations:
        p = tf.add_paragraph()
        p.text = f"✓ {obs}"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_TEXT
        p.space_after = Pt(2)

    return slide


# ==================== CREATE SLIDES ====================

# SLIDE 1: TITLE
add_title_slide(
    prs,
    "Brain Tumor Diagnosis Assistant",
    "Group 8\n\nTarik Bilgin Demirci – Technical Lead (Student 1)\nUmut Turklay – Figures, Tables & Presentation (Student 2)\nBerk Kahraman – Report & Storytelling (Student 3)\n\nPattern Recognition | Winter 2025"
)

# SLIDE 2: PROBLEM STATEMENT
add_content_slide(
    prs,
    "Problem Statement",
    [
        "Brain tumors are life-threatening conditions requiring accurate and timely diagnosis",
        "Manual MRI analysis is time-consuming and prone to inter-observer variability",
        "Need for automated, reliable, and interpretable diagnostic systems",
        "Challenge: Achieve high accuracy while maintaining clinical trust",
        "Goal: Develop a hybrid AI system combining deep learning with clinical knowledge"
    ],
    slide_num=2
)

# SLIDE 3: RESEARCH QUESTIONS
add_content_slide(
    prs,
    "Research Questions",
    [
        "RQ1: How effectively can CNNs classify brain tumor types from MRI scans?",
        "RQ2: Can meta-learning with shape descriptors improve classification accuracy?",
        "RQ3: How can rule-based clinical knowledge enhance diagnostic reliability?",
        "RQ4: What is the impact of explainability techniques (Grad-CAM) on interpretability?",
        "RQ5: How do different architectures (ResNet, EfficientNet, DenseNet) compare?"
    ],
    slide_num=3
)

# SLIDE 4: LITERATURE REVIEW
add_table_slide(
    prs,
    "Literature Review / Related Work",
    ["Author", "Year", "Method", "Dataset", "Accuracy", "Limitation"],
    [
        ["Deepak & Ameer", "2019", "GoogleNet + SVM", "CE-MRI (3064)", "98.0%", "No explainability"],
        ["Afshar et al.", "2019", "CapsNet", "CE-MRI (3064)", "90.9%", "Limited tumor types"],
        ["Badza & Barjaktarovic", "2020", "CNN (Custom)", "Figshare (3064)", "96.6%", "No clinical integration"],
        ["Swati et al.", "2019", "VGG-19 Transfer", "CE-MRI (3064)", "94.8%", "Single architecture"],
        ["Our Work", "2025", "Xception+Meta+Rules", "Kaggle (7023)", "99.5%", "Hybrid approach"]
    ],
    "Our work: Combining CNN with meta-learning and rule-based clinical integration for improved reliability.",
    slide_num=4
)

# SLIDE 5: DATASET
add_figure_slide(
    prs,
    "Dataset Overview",
    png_files.get("RQ1_Fig3"),
    [
        "Source: Kaggle Brain Tumor MRI Dataset | 4 Classes: Glioma, Meningioma, Pituitary, No Tumor",
        "Total: 5,712 training | 655 validation | 656 test images | Size: 299×299 RGB"
    ],
    slide_num=5
)

# SLIDE 6: NETWORK ARCHITECTURE - XCEPTION (fullsize)
add_figure_slide(
    prs,
    "Network Architecture: Xception (Primary Model)",
    f"{png_dir}/Architecture_Xception.png",
    slide_num=6,
    fullsize=True
)

# SLIDE 7: NETWORK ARCHITECTURE - COMPARISON (fullsize)
add_figure_slide(
    prs,
    "Network Architecture: Comparison Models (RQ5)",
    f"{png_dir}/Architecture_Comparison.png",
    slide_num=7,
    fullsize=True
)

# SLIDE 8: METHODOLOGY (fullsize)
add_figure_slide(
    prs,
    "Methodology: Hybrid Classification Pipeline",
    f"{png_dir}/Methodology_Pipeline.png",
    slide_num=8,
    fullsize=True
)

# SLIDE 9: RESULTS RQ1
add_rq_result_slide(
    prs,
    rq_num=1,
    rq_title="CNN Classification Effectiveness",
    image_path=png_files.get("RQ1_Fig4"),
    table_headers=["Class", "Precision", "Recall", "F1", "Support"],
    table_rows=[
        ["Glioma", "0.99", "0.99", "0.99", "150"],
        ["Meningioma", "0.99", "0.99", "0.99", "153"],
        ["No Tumor", "1.00", "1.00", "1.00", "203"],
        ["Pituitary", "1.00", "0.99", "1.00", "150"],
        ["Overall", "0.99", "0.99", "0.99", "656"]
    ],
    observations=[
        "Xception achieves 99.39% test accuracy on 4-class brain tumor classification",
        "Perfect recall (100%) on No Tumor class - critical for avoiding false negatives",
        "Balanced performance across all tumor types with F1 ≥ 0.99"
    ],
    slide_num=9
)

# SLIDE 10: RESULTS RQ2
add_rq_result_slide(
    prs,
    rq_num=2,
    rq_title="Meta-Learning with Shape Features",
    image_path=png_files.get("RQ2_Fig1"),
    table_headers=["Model", "Accuracy", "Improvement"],
    table_rows=[
        ["CNN Only (Xception)", "99.39%", "Baseline"],
        ["Meta-Learner (CNN + Shape)", "99.54%", "+0.15%"]
    ],
    observations=[
        "Shape features (area, perimeter, circularity, solidity, irregularity) provide complementary information",
        "Meta-learner (Logistic Regression) successfully combines CNN outputs with shape descriptors",
        "Modest but consistent improvement demonstrates value of hybrid approach"
    ],
    slide_num=10
)

# SLIDE 11: RESULTS RQ3
add_rq_result_slide(
    prs,
    rq_num=3,
    rq_title="Rule-Based Clinical Integration",
    image_path=png_files.get("RQ3_Fig1"),
    table_headers=["Metric", "Value"],
    table_rows=[
        ["Confidence Threshold", "0.80"],
        ["Irregularity Threshold", "26.81 (90th pctl)"],
        ["ACCEPT Rate", "89.18% (585)"],
        ["REFER Rate", "10.82% (71)"],
        ["Accuracy on ACCEPT", "99.66%"]
    ],
    observations=[
        "Rule engine flags uncertain cases based on confidence and tumor irregularity",
        "89% of cases automatically accepted with higher accuracy (99.66%)",
        "Remaining 11% referred to specialists - practical clinical workflow"
    ],
    slide_num=11
)

# SLIDE 12: RESULTS RQ4
add_rq_result_slide(
    prs,
    rq_num=4,
    rq_title="Explainability with Grad-CAM",
    image_path=png_files.get("RQ4_Fig1"),
    table_headers=["Aspect", "Description"],
    table_rows=[
        ["Method", "Gradient-weighted Class Activation Mapping"],
        ["Target Layer", "block14_sepconv2_act (Xception)"],
        ["Output", "Heatmap highlighting influential regions"],
        ["Purpose", "Enable clinician trust and understanding"]
    ],
    observations=[
        "Grad-CAM successfully highlights tumor regions that influence predictions",
        "Provides visual explanation for each classification decision",
        "Critical for clinical adoption and regulatory compliance (explainable AI)"
    ],
    slide_num=12
)

# SLIDE 13: RESULTS RQ5
add_rq_result_slide(
    prs,
    rq_num=5,
    rq_title="Architecture Comparison",
    image_path=png_files.get("RQ5_Fig1"),
    table_headers=["Architecture", "Accuracy", "Params", "Time (s)"],
    table_rows=[
        ["Xception", "97.87%", "21.1M", "1,390"],
        ["DenseNet121", "95.88%", "7.2M", "2,431"],
        ["ResNet50", "32.93%", "23.9M", "2,416"],
        ["EfficientNetB0", "30.95%", "4.2M", "1,999"]
    ],
    observations=[
        "Xception significantly outperforms all other architectures for this task",
        "DenseNet121 shows reasonable performance but lower than Xception",
        "ResNet50 and EfficientNetB0 fail to generalize - need different fine-tuning"
    ],
    slide_num=13
)

# SLIDE 14: CONCLUSION
add_content_slide(
    prs,
    "Conclusion",
    [
        "Key Findings:",
        "   • Xception CNN achieves 99.39% accuracy on brain tumor classification",
        "   • Meta-learning with shape features improves accuracy to 99.54%",
        "   • Rule-based engine enables 99.66% accuracy on 89% auto-accepted cases",
        "   • Grad-CAM provides interpretable visualizations for clinical trust",
        "",
        "Limitations:",
        "   • Single dataset - needs multi-center validation",
        "   • Limited to 4 tumor types",
        "",
        "Future Work:",
        "   • Extend to more tumor types and 3D MRI volumes",
        "   • Clinical trials for real-world deployment validation"
    ],
    slide_num=14
)

# Save
output_path = "/Users/umutturklay/dev/UE/pattern_recog/Group8_BrainTumorDiagnosis_Part3_Presentation.pptx"
prs.save(output_path)
print(f"\nPresentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
print("Design: Medical/Health theme with blue-teal color scheme")
